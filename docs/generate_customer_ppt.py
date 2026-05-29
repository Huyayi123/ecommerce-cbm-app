from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt


FONT = "Microsoft YaHei"
OUT = Path(__file__).with_name("customer-presentation-erp-fixed.pptx")

COLORS = {
    "green": RGBColor(19, 111, 99),
    "dark": RGBColor(22, 32, 42),
    "muted": RGBColor(95, 109, 120),
    "line": RGBColor(203, 214, 223),
    "pale": RGBColor(238, 248, 245),
    "bluepale": RGBColor(238, 245, 248),
    "white": RGBColor(255, 255, 255),
}


def set_run_font(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        xml = f'<a:{tag} {nsdecls("a")} typeface="{FONT}"/>'
        node = parse_xml(xml)
        old = rpr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if old is not None:
            rpr.remove(old)
        rpr.append(node)


def set_text(paragraph, text, size, color, bold=False, align=None):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, size, color, bold)
    if align is not None:
        paragraph.alignment = align


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    set_text(box.text_frame.paragraphs[0], title, 28, COLORS["dark"], True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(11.8), Inches(0.45))
        set_text(sub.text_frame.paragraphs[0], subtitle, 13, COLORS["muted"])


def add_footer(slide):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.22), Inches(12), Inches(0.2))
    set_text(foot.text_frame.paragraphs[0], "电商采购装柜轻型 ERP | 客户展示版", 8, COLORS["muted"])


def add_box(slide, x, y, w, h, text, fill="white", size=13, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for index, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        set_text(paragraph, line, size, COLORS["white"] if fill == "green" else COLORS["dark"], bold, PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = COLORS["green"]
    connector.line.width = Pt(1.5)
    connector.line.end_arrowhead = True


def add_bullets(slide, x, y, w, h, items, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        set_text(paragraph, item, size, COLORS["dark"])
        paragraph.space_after = Pt(8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "电商采购装柜轻型 ERP", "从 Takealot 销量预测到采购任务、装柜计算、海运在途跟踪的一体化系统")
    add_box(slide, 0.85, 1.75, 3.4, 1.1, "自动同步\nTakealot 销量与库存", "bluepale", 17, True)
    add_box(slide, 4.95, 1.75, 3.4, 1.1, "智能生成\n采购建议", "pale", 17, True)
    add_box(slide, 9.05, 1.75, 3.4, 1.1, "跟踪采购\n与海运在途", "bluepale", 17, True)
    add_arrow(slide, 4.25, 2.3, 4.95, 2.3)
    add_arrow(slide, 8.35, 2.3, 9.05, 2.3)
    add_bullets(slide, 1.25, 3.7, 10.8, 1.8, [
        "目标：减少人工表格整理、避免漏采错采、控制整柜 68-70 CBM",
        "适用：多店铺、多采购人、海运周期长、需要持续补货预测的电商团队",
        "结果：采购建议可调整，采购任务可分配，在途库存可追踪，验货单可导出",
    ], 16)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "系统总流程", "数据进入系统，经过预测、装柜、采购确认，最终形成在途库存")
    steps = [
        ("Takealot API\n销量 / 库存", 0.45),
        ("月销量\n采购建议", 2.45),
        ("采购池\n按店铺汇总", 4.45),
        ("装柜计算\n68-70 CBM", 6.45),
        ("生成\n采购任务", 8.45),
        ("我的\n采购订单", 10.45),
    ]
    prev = None
    for text, x in steps:
        add_box(slide, x, 1.65, 1.55, 0.85, text, "pale", 12, True)
        if prev is not None:
            add_arrow(slide, prev + 1.55, 2.07, x, 2.07)
        prev = x
    add_box(slide, 4.45, 3.6, 1.55, 0.85, "SKU 资料库\n价格 / CBM / 采购人", "bluepale", 11, True)
    add_arrow(slide, 5.22, 3.6, 3.22, 2.5)
    add_box(slide, 10.45, 3.6, 1.55, 0.85, "采购 / 海运\n在途库存", "green", 11, True)
    add_arrow(slide, 11.22, 2.5, 11.22, 3.6)
    add_arrow(slide, 10.45, 4.05, 1.25, 4.05)
    add_box(slide, 0.45, 3.6, 1.55, 0.85, "参与下一轮\n采购预测", "bluepale", 11, True)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "五大核心页面职责", "每个页面只负责一个清晰环节，避免采购计划、采购执行和物流数据混在一起")
    modules = [
        ("SKU 资料库", "商品主数据：SKU、厂家、价格、CBM、采购人"),
        ("月销量采购建议", "同步 API，计算建议采购数量，可人工修改"),
        ("装柜计算", "控制柜容，调整数量，生成采购任务"),
        ("我的采购订单", "采购人查看自己的任务并确认采购"),
        ("采购 / 在途库存", "维护已确认采购、物流数据、到货状态、导出验货单"),
    ]
    for index, (title, desc) in enumerate(modules):
        y = 1.35 + index * 0.85
        add_box(slide, 0.8, y, 2.2, 0.55, title, "green", 13, True)
        add_box(slide, 3.25, y, 8.9, 0.55, desc, "white", 13, False)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "采购建议计算逻辑", "把销售、库存、在途库存统一成一套可解释的订货公式")
    add_box(slide, 0.75, 1.35, 11.85, 0.85, "建议采购数量 = 月销量 × 备货月数 - 南非本地库存 - 官方仓库存 - 送仓路上库存 - 海运在途库存", "pale", 17, True)
    add_bullets(slide, 1.0, 2.55, 5.6, 2.2, [
        "月销量 > 50：备货 4 个月",
        "月销量 21-50：备货 3 个月",
        "月销量 20 及以下：备货 2 个月",
        "建议采购数量可以人工调整",
    ], 15)
    add_bullets(slide, 7.0, 2.55, 5.3, 2.2, [
        "SKU 越大越新，用于新品判断",
        "Bestby / Arfast / Aicom 有不同新品倍率",
        "新品预测在备注里展示，不阻止采购",
        "未录入 SKU 资料才标红提示",
    ], 15)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "采购执行 SOP", "从系统预测到采购确认，再到在途库存管理")
    ops = ["维护 SKU 资料", "同步 Takealot 数据", "生成并调整采购建议", "加入采购池", "装柜计算 68-70 CBM", "生成采购任务", "采购人确认订单", "进入采购 / 海运在途库存"]
    for index, op in enumerate(ops):
        x = 0.55 + (index % 4) * 3.15
        y = 1.45 + (index // 4) * 2.0
        add_box(slide, x, y, 2.55, 0.75, f"{index + 1}. {op}", "pale" if index < 4 else "bluepale", 13, True)
        if index % 4 != 3:
            add_arrow(slide, x + 2.55, y + 0.38, x + 3.15, y + 0.38)
    add_arrow(slide, 10.0, 2.2, 1.83, 3.45)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "角色权限与协作方式", "内部多人在线使用，各角色看到和维护自己负责的数据")
    roles = [
        ("admin 管理员", ["维护 SKU", "查看和编辑全部数据", "管理采购和在途库存"]),
        ("buyer 采购员", ["查看自己的采购订单", "确认数量、价格、状态", "补充备注和物流数据"]),
        ("viewer 查看人员", ["只读查看", "不修改数据", "适合管理层或协作人员"]),
    ]
    for index, (role, items) in enumerate(roles):
        x = 0.75 + index * 4.15
        add_box(slide, x, 1.45, 3.4, 0.7, role, "green", 15, True)
        add_bullets(slide, x + 0.15, 2.35, 3.1, 2.0, items, 14)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "客户演示顺序", "建议按真实业务路径录屏或现场演示，控制在 3-5 分钟")
    add_bullets(slide, 0.8, 1.25, 11.8, 5.2, [
        "1. 登录系统，说明这是内部多人在线工作台",
        "2. 展示 SKU 资料库，说明商品主数据来源",
        "3. 打开月销量采购建议，展示 Takealot API 同步和新品预测",
        "4. 修改建议采购数量，加入本轮采购池",
        "5. 发送到装柜计算，调整到 68-70 CBM",
        "6. 生成采购任务，展示采购人顶部待采购提醒",
        "7. 采购人在我的采购订单确认数量、价格和状态",
        "8. 进入采购 / 海运在途库存，补充物流数据并导出验货单",
    ], 15)
    add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "客户价值总结", "把分散表格变成可预测、可分配、可追踪的采购闭环")
    add_bullets(slide, 1.0, 1.35, 11.0, 4.7, [
        "效率：自动同步销量和库存，减少人工整理报表",
        "准确：结合库存和海运在途，减少重复采购和漏采",
        "装柜：采购建议进入装柜计算，控制 68-70 CBM",
        "协作：采购任务按采购人分配，登录后可看到待采购提醒",
        "追踪：采购确认后进入在途库存，物流数据和到货状态可维护",
        "沉淀：历史采购和在途数据参与下一轮采购预测",
    ], 17)
    add_footer(slide)

    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
